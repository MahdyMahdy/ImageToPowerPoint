import cv2
import numpy as np
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

class Shape():
    def __init__(self,center,edges):
        self.center = center
        self.edges = edges
        self.nei = []
        self.reel = None
        pass

class Line():
    def __init__(self,a,b,xs):
        self.a = a
        self.b=b
        self.xs = xs
        self.nei = []

def shapes_bw(orig):
    image = orig.copy()
    gray=cv2.cvtColor(image,cv2.COLOR_BGR2GRAY)
    (thresh, image_bw) = cv2.threshold(gray, 128, 255, cv2.THRESH_BINARY_INV)
    kernel = np.ones((3,3), np.uint8)
    img_dilate = cv2.dilate(image_bw, kernel, iterations=15)
    h, w = img_dilate.shape[:2]
    img_fool = img_dilate.copy()
    mask = np.zeros((h+2, w+2), np.uint8)
    cv2.floodFill(img_fool, mask, (0,0), 255)
    inv = cv2.bitwise_not(img_fool)
    im_out = img_dilate | inv
    kernel = np.ones((5,5), np.uint8)
    img_ero = cv2.erode(im_out, kernel, iterations=10)
    (thresh, result) = cv2.threshold(img_ero, 128, 255, cv2.THRESH_BINARY_INV)
    return result

def detect_shapes(orig):
    shapes = []
    image = orig.copy()
    step1 = shapes_bw(image)
    thresh = cv2.adaptiveThreshold(step1,255,cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY_INV,51,9)
    cnts = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    cnts = cnts[0] if len(cnts) == 2 else cnts[1]
    for c in cnts:
        cv2.drawContours(thresh, [c], -1, (255,255,255), -1)
    kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (9,9))
    opening = cv2.morphologyEx(thresh, cv2.MORPH_OPEN, kernel, iterations=4)
    cnts = cv2.findContours(opening, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    cnts = cnts[0] if len(cnts) == 2 else cnts[1]
    area_treshold = 4000
    for c in cnts:
        if cv2.contourArea(c) > area_treshold :
            approx = cv2.approxPolyDP(c, 0.01*cv2.arcLength(c, True), True)
            if len(approx) >=3:
                edges = 0
                for i,p1 in enumerate(c[1:]):
                    p0=c[i]
                    if np.linalg.norm(p0-p1)>15:
                        edges+=1
                M = cv2.moments(c)
                if M['m00'] != 0:
                    cx = int(M['m10']/M['m00'])
                    cy = int(M['m01']/M['m00'])
                shape = Shape((cx/image.shape[1],cy/image.shape[0]),edges)
                shapes.append(shape)
                cv2.drawContours(orig, [approx], 0, (0,0,255), 4)
    return shapes

def dist(x1,y1,x2,y2):
    z = np.power(x1-x2,2)+np.power(y1-y2,2)
    return np.sqrt(z)

def detect_neigh(shapes,lines,img_size):
    if lines == None:
        return
    for line in lines:
        a = line.a
        b = line.b
        i=0
        for x in line.xs:
            d = float("inf")
            y = a*x+b
            for shape in shapes:
                cx = shape.center[0]*img_size[1]
                cy = shape.center[1]*img_size[0]
                z=dist(x,y,cx,cy)
                if z<d:
                    d = z
                    if len(line.nei)==i:
                        line.nei.append(shape)
                    else:
                        line.nei[i]=shape
            i+=1
    for line in lines:
        n1 =line.nei[0]
        n2 =line.nei[1]
        side = 1
        if n1.center[0]<n2.center[0]:
            side = 3
        n1.nei.append((n2,side))
        if side == 3:
            side = 1
        else:
            side = 3
        n2.nei.append((n1,side))

def detectLines(image):
    step1 = getLines(image)
    edges = cv2.Canny(step1, 50, 150)
    lines = cv2.HoughLinesP(edges,3,np.pi/180,50,minLineLength=20,maxLineGap=20)
    if lines is None:
        lines = []
    slopes = []
    bs = []
    maxs = []
    mins = []
    for line in lines:
        for x1,y1,x2,y2 in line:
            cv2.line(image,(x1,y1),(x2,y2),(255,0,0),5)
            a = (y2-y1)/(x2-x1)
            same = False
            for i,slope in enumerate(slopes):
                if np.abs(a-slope)<1 and (np.abs(maxs[i]-max(x1,x2))<30 or 
                 np.abs(mins[i]-min(x1,x2))<30 or 
                 np.abs(mins[i]-max(x1,x2))<30 or 
                 np.abs(maxs[i]-min(x1,x2))<30):
                    slope = (a+slope)/2
                    maxs[i] = max(maxs[i],x1,x2)
                    mins[i] = min(mins[i],x1,x2)
                    same = True
            if not same or len(slopes)==0:
                slopes.append(a)
                b = y2 - a*x2
                bs.append(b)
                maxs.append(max(x1,x2))
                mins.append(min(x1,x2))
    lines = []
    for i in range(0,len(bs)):
        line = Line(slopes[i],bs[i],(maxs[i],mins[i]))
        lines.append(line)
    return lines

def getLines(orig):
    image = shapes_bw(orig)
    gray=cv2.cvtColor(orig,cv2.COLOR_BGR2GRAY)
    (thresh, orig_bw) = cv2.threshold(gray, 128, 255, cv2.THRESH_BINARY_INV)
    kernel = np.ones((5,5), np.uint8)
    img_ero = cv2.erode(image, kernel, iterations=10)
    return orig_bw & img_ero

def create_pp(shapes_to_insert):
    x = Presentation()
    layout = x.slide_layouts[6] 
    first_slide = x.slides.add_slide(layout)
    done = []
    for shape in shapes_to_insert:
        if not shape in done:
            done.extend(draw_shapes(shape,first_slide))
            done=list(set(done))
    x.slide_width = Inches(16)
    x.slide_height = Inches(9)
    x.save("presentation.pptx")

def draw_shapes(shape,slide):
    q = []
    done = []
    q.append(shape)
    for shape in q:
        left = Inches(shape.center[0]*16)
        top = Inches(shape.center[1]*9)
        width = height = Inches(1.0)
        if shape.edges == 3:
            kind = MSO_SHAPE.ISOSCELES_TRIANGLE
        elif shape.edges == 4:
            kind = MSO_SHAPE.RECTANGLE
        elif shape.edges == 5:
            kind = MSO_SHAPE.PENTAGON
        elif shape.edges == 6:
            kind = MSO_SHAPE.HEXAGON
        elif shape.edges == 7:
            kind = MSO_SHAPE.HEPTAGON
        elif shape.edges == 8:
            kind = MSO_SHAPE.OCTAGON
        elif shape.edges == 9:
            kind = MSO_SHAPE.PENTAGON
        elif shape.edges == 10:
            kind = MSO_SHAPE.DECAGON
        else:
            kind = MSO_SHAPE.OVAL
        shape.reel =slide.shapes.add_shape(kind, left, top, width, height)
        done.append(shape)
        for nei in shape.nei:
            if not nei[0] in done:
                q.append(nei[0])
            else:
                line = slide.shapes.add_connector(1,Inches(0.5*9),Inches(0.5*9),Inches(0.7*9),Inches(0.7*9))
                side = nei[1]
                if side == 3:
                    side = 1
                else:
                    side = 3
                line.begin_connect(shape.reel,nei[1])
                line.end_connect(nei[0].reel,side)
    return done

